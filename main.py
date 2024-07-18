from fastapi import FastAPI, HTTPException, Form
from pydantic import BaseModel
import requests
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
import boto3
from datetime import datetime
import re
import os
from dotenv import load_dotenv

load_dotenv()
aws_access_key_id = os.getenv("AWS_ACCESS_KEY_ID")
aws_secret_access_key = os.getenv("AWS_SECRET_ACCESS_KEY")
aws_region = os.getenv("AWS_REGION")
aws_s3_bucket_name = os.getenv("AWS_S3_BUCKET_NAME")

s3 = boto3.client('s3',
                  aws_access_key_id=aws_access_key_id,
                  aws_secret_access_key=aws_secret_access_key,
                  region_name=aws_region)

app = FastAPI()

class ExtractImagesRequest(BaseModel):
    url: str
    client : str

@app.post("/extract_image_and_location_from_pptx")
async def extract_and_process(request: ExtractImagesRequest):
    try:
        if request.client == 'kaushik':
            result = extract_locations_and_image_kaushik(request.url)
            return result
        elif request.client == 'mantra':
            result = extract_location_and_image_mantra(request.url)
            return result
        elif request.client == 'chitra':
            result = extract_location_and_image_chitra(request.url)
            return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def extract_locations_and_image_kaushik(pptx_s3_url, crop_percentage=7):
    response = requests.get(pptx_s3_url)
    pptx_bytes = BytesIO(response.content)
    pptx_bytes.seek(0)
    presentation = Presentation(pptx_bytes)
    result_list = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_title = shape.text_frame.text.strip().replace(" ", "_")
                break
        if not slide_title:
            slide_title = f"slide_{slide_number}"
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                width, height = image.size
                crop_height = int(height * crop_percentage / 100)
                crop_box = (0, height - crop_height, width, height)
                cropped_image = image.crop(crop_box)
                res = pytesseract.image_to_string(cropped_image)
                res = res.strip().lower()
                text = res.replace(' ','')
                pattern = r'^[a-zA-Z0-9()]+-(.*?)(\d{2}x\d{2})'
                # Search for the pattern in the input string
                match = re.search(pattern, text)

                # Extract the desired part of the string if a match is found
                if match:
                    text = match.group(1).strip()  # .strip() removes any leading/trailing whitespace
                    text = re.sub(r'[^\w\s]', '', text)
                    
                remaining_image = image.crop((0, 0, width, height - crop_height))
                remaining_image_bytes = BytesIO()
                remaining_image.save(remaining_image_bytes, format=image.format)
                remaining_image_bytes.seek(0)
                folder_structure = 'prod/ldoc/inventoryManagement'
                timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
                s3_key = f"{folder_structure}/{slide_title}_shape_{shape_number}_{timestamp}.{image.format.lower()}"
                s3.upload_fileobj(remaining_image_bytes, aws_s3_bucket_name, s3_key,ExtraArgs={'ACL': 'public-read'})
                s3_url = f"https://{aws_s3_bucket_name}.s3.{aws_region}.amazonaws.com/{s3_key}"
                result_list.append({text: s3_url})
    return result_list


def extract_location_and_image_mantra(ppt_url):
    response = requests.get(ppt_url)
    pptx_bytes = BytesIO(response.content)
    pptx_bytes.seek(0)
    prs = Presentation(pptx_bytes)
    result_list=[]

    for slide_index in range(len(prs.slides)):
        slide_title = None
        if not slide_title:
            slide_title = f"slide_{slide_index}"
        slide = prs.slides[slide_index]
        slide_info={}
        has_image = False
        has_table = False
        for shape_number,shape in enumerate(slide.shapes,start=1):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                if shape.image:
                    image_part = shape.image
                    image_stream = BytesIO(image_part.blob)
                    image= Image.open(image_stream)
                    image_bytes = BytesIO()
                    image.save(image_bytes , format=image.format)
                    image_bytes.seek(0)
                    folder_structure = 'prod/ldoc/inventoryManagement'
                    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
                    s3_key = f"{folder_structure}/{slide_title}_shape_{shape_number}_{timestamp}.{image.format.lower()}"
                    s3.upload_fileobj(image_bytes, aws_s3_bucket_name, s3_key,ExtraArgs={'ACL': 'public-read'})
                    s3_url = f"https://{aws_s3_bucket_name}.s3.{aws_region}.amazonaws.com/{s3_key}"
                    slide_info['url'] = s3_url

            if shape.shape_type == 19:
                has_table = True
                table = shape.table
                if table.rows and table.columns:
                    text = table.cell(0,0).text.strip()
                    if text[0].isdigit():
                        text = text[2:].strip()
                    stripped_text = text.strip().lower()
                    text = stripped_text.replace(' ','')
                    text = text.split("-")[2:-2][0]
                    result = re.sub(r'[^\w\s]', '', text)
                    slide_info['location'] = result
            
        if has_image and has_table:
            result_list.append({slide_info['location']:slide_info['url']})
    return result_list


def extract_location_and_image_mantra(ppt_url):
    response = requests.get(ppt_url)
    pptx_bytes = BytesIO(response.content)
    pptx_bytes.seek(0)
    prs = Presentation(pptx_bytes)
    result_list=[]

    for slide_index in range(1,len(prs.slides) -1):
        slide_title = None
        if not slide_title:
            slide_title = f"slide_{slide_index}"
        slide = prs.slides[slide_index]
        slide_info={}
        has_image = False
        has_table = False
        for shape_number,shape in enumerate(slide.shapes,start=1):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                if shape.image:
                    image_part = shape.image
                    image_stream = BytesIO(image_part.blob)
                    image= Image.open(image_stream)
                    image_bytes = BytesIO()
                    image.save(image_bytes , format=image.format)
                    image_bytes.seek(0)
                    folder_structure = 'prod/ldoc/inventoryManagement'
                    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
                    s3_key = f"{folder_structure}/{slide_title}_shape_{shape_number}_{timestamp}.{image.format.lower()}"
                    s3.upload_fileobj(image_bytes, aws_s3_bucket_name, s3_key,ExtraArgs={'ACL': 'public-read'})
                    s3_url = f"https://{aws_s3_bucket_name}.s3.{aws_region}.amazonaws.com/{s3_key}"
                    slide_info['url'] = s3_url

            if shape.shape_type == 19:
                has_table = True
                table = shape.table
                if table.rows and table.columns:
                    text = table.cell(0,0).text.strip()
                    if text[0].isdigit():
                        text = text[2:].strip()
                    slide_info['location'] = '*' + text + '*'
            
        if has_image and has_table:
            result_list.append({slide_info['location']:slide_info['url']})
    return result_list

def extract_location(text):
    parts = text.split('-')
    
    # Check if there's at least one hyphen in the text
    if len(parts) > 1:
        # Take the last part and strip any leading/trailing whitespace
        location = parts[-1].strip()
        lowered_text = location.lower()
        text = lowered_text.replace(' ','')
        result = re.sub(r'[^\w\s]', '', text)
        return result
    else:
        return text


def extract_location_and_image_chitra(ppt_url):
    response = requests.get(ppt_url)
    pptx_bytes = BytesIO(response.content)
    pptx_bytes.seek(0)
    prs = Presentation(pptx_bytes)
    result_list = []

    for slide_index, slide in enumerate(prs.slides):
        slide_title = f"slide_{slide_index}"
        slide_info = {}
        has_image = False
        has_table = False

        for shape_number, shape in enumerate(slide.shapes, start=1):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                try:
                    # Method 1: Try to get image through shape.image
                    image_part = shape.image
                    image_ext = image_part.ext
                    image_data = image_part.blob
                except AttributeError:
                    try:
                        # Method 2: Try to get image through slide.part.rels
                        blip = shape._element.xpath('.//a:blip')[0]
                        rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        image_part = slide.part.rels[rId].target_part
                        
                        image_ext = 'png'  # Default to png if no extension is found
                        
                        image_data = image_part.blob
                    except Exception as e:
                        print(f"Failed to extract image from shape {shape_number} on slide {slide_index}: {str(e)}")
                        continue

                image_stream = BytesIO(image_data)
                image = Image.open(image_stream)
                image_bytes = BytesIO()
                image.save(image_bytes, format=image.format if image.format else 'PNG')
                image_bytes.seek(0)

                folder_structure = 'prod/ldoc/inventoryManagement'
                timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
                s3_key = f"{folder_structure}/{slide_title}_shape_{shape_number}_{timestamp}.{image_ext}"

                try:
                    s3.upload_fileobj(image_bytes, aws_s3_bucket_name, s3_key, ExtraArgs={'ACL': 'public-read'})
                    s3_url = f"https://{aws_s3_bucket_name}.s3.{aws_region}.amazonaws.com/{s3_key}"
                    slide_info['url'] = s3_url
                except Exception as e:
                    print(f"Error uploading to S3: {str(e)}")
            
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                has_table = True
                table = shape.table
                if table.rows and table.columns:
                    text = table.cell(0, 0).text.strip()
                    location = extract_location(text)
                    if location:
                        slide_info['location'] = location

        if has_image and has_table and 'location' in slide_info and 'url' in slide_info:
            result_list.append({slide_info['location']: slide_info['url']})

    return result_list



if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
